# make_calendar.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import calendar
import sys
import math

# Color themes (RGB values for header backgrounds)
COLOR_THEMES = {
    "1": {"name": "Blue", "rgb": (70, 130, 180)},
    "2": {"name": "Navy", "rgb": (31, 73, 125)},
    "3": {"name": "Teal", "rgb": (0, 128, 128)},
    "4": {"name": "Purple", "rgb": (147, 112, 219)},
    "5": {"name": "Green", "rgb": (34, 139, 34)},
    "6": {"name": "Orange", "rgb": (255, 140, 0)},
    "7": {"name": "Red", "rgb": (205, 92, 92)},
    "8": {"name": "Pink", "rgb": (219, 112, 147)},
    "9": {"name": "Gray", "rgb": (128, 128, 128)},
    "10": {"name": "Dark Gray", "rgb": (64, 64, 64)},
    "11": {"name": "Black", "rgb": (0, 0, 0)},
    "12": {"name": "White", "rgb": (255, 255, 255)},
}

def build_year_calendar(year: int, filename: str, theme_rgb: tuple):
    prs = Presentation()
    
    # Set slide size to widescreen 16:9 (standard PowerPoint widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Slide size (EMU units)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Layout configuration
    cols = 4  # months per row
    rows = 3  # rows of months
    
    # Each calendar is 3.04" wide x 2" tall with 0.2" gutters between columns
    calendar_width = Inches(3.04)
    calendar_height = Inches(2.0)
    h_gutter = Inches(0.2)
    v_gutter = Inches(0.3)
    
    # Calculate margins to center the calendars
    total_width = cols * calendar_width + (cols - 1) * h_gutter
    total_height = rows * calendar_height + (rows - 1) * v_gutter
    
    margin_left = (slide_width - total_width) / 2
    margin_right = margin_left
    margin_top = (slide_height - total_height) / 2
    margin_bottom = margin_top
    
    title_height = Inches(0.3)

    month_names = [
        "January", "February", "March", "April",
        "May", "June", "July", "August",
        "September", "October", "November", "December"
    ]
    weekday_labels = ["S", "M", "T", "W", "T", "F", "S"]

    for i, name in enumerate(month_names):
        month_num = i + 1
        col = i % cols
        row = i // cols

        # Calculate positions - ensure they're integers
        left = int(margin_left + col * (calendar_width + h_gutter))
        top = int(margin_top + row * (calendar_height + v_gutter))

        # ----- Month title -----
        title_box = slide.shapes.add_textbox(left, top, calendar_width, title_height)
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{name} {year}"
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Calculate rows needed for this specific month
        first_weekday, num_days = calendar.monthrange(year, month_num)
        start_col = (first_weekday + 1) % 7  # Adjust to Sunday=0
        
        # Calculate how many weeks (rows) we need
        days_in_first_week = 7 - start_col
        remaining_days = num_days - days_in_first_week
        additional_weeks = math.ceil(remaining_days / 7) if remaining_days > 0 else 0
        weeks_needed = 1 + additional_weeks
        
        # Total rows = 1 header + weeks_needed
        table_rows = 1 + weeks_needed
        table_cols = 7

        # Table dimensions - should be 3.04" wide 
        table_width = calendar_width  # Keep as EMU, don't convert to int yet
        
        # Available height for table (after title and small gap)
        table_gap = Inches(0.02)
        available_table_height = calendar_height - title_height - table_gap
        
        # Use the full available height for the table
        table_height = available_table_height
        table_top = int(top + title_height + table_gap)

        table_shape = slide.shapes.add_table(
            table_rows, table_cols, left, table_top, int(table_width), int(table_height)
        )
        table = table_shape.table

        # Header row with theme color
        for c in range(table_cols):
            cell = table.cell(0, c)
            cell.text = weekday_labels[c]
            
            # Apply theme color to header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*theme_rgb)
            
            # Set font color based on background brightness
            r, g, b = theme_rgb
            brightness = (r * 299 + g * 587 + b * 114) / 1000
            font_color = RGBColor(255, 255, 255) if brightness < 128 else RGBColor(0, 0, 0)
            
            cell.text_frame.paragraphs[0].font.name = "Arial"
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = font_color
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Fill in the days
        day = 1
        for r in range(1, table_rows):
            for c in range(table_cols):
                # Only add day if we're past the start column on first row
                # or if we haven't exceeded num_days
                if (r == 1 and c >= start_col and day <= num_days) or \
                   (r > 1 and day <= num_days):
                    cell = table.cell(r, c)
                    cell.text = str(day)
                    cell.text_frame.paragraphs[0].font.name = "Arial"
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    day += 1
                else:
                    # Empty cell
                    cell = table.cell(r, c)
                    cell.text = ""
                    cell.text_frame.paragraphs[0].font.name = "Arial"
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(filename)
    print(f"Saved calendar for {year} to {filename}")


if __name__ == "__main__":
    # Usage: python make_calendar.py 2026
    if len(sys.argv) >= 2:
        year = int(sys.argv[1])
    else:
        year = int(input("Enter year (e.g. 2026): "))
    
    # Display color theme options
    print("\nChoose a color theme:")
    print("Cool Tones:")
    print("  1. Blue (Sky Blue)")
    print("  2. Navy (Professional Blue)")
    print("  3. Teal")
    print("  4. Purple")
    print("\nWarm Tones:")
    print("  5. Green")
    print("  6. Orange")
    print("  7. Red")
    print("  8. Pink")
    print("\nNeutral/Monochrome:")
    print("  9. Gray")
    print("  10. Dark Gray")
    print("  11. Black")
    print("  12. White")
    
    theme_choice = input("\nEnter theme number (1-12, default is 1-Blue): ").strip()
    if not theme_choice:
        theme_choice = "1"
    
    if theme_choice not in COLOR_THEMES:
        print(f"Invalid choice, using Blue theme")
        theme_choice = "1"
    
    theme = COLOR_THEMES[theme_choice]
    print(f"Using {theme['name']} theme")

    output = f"Calendar_{year}.pptx"
    build_year_calendar(year, output, theme['rgb'])
